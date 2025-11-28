import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def hex_to_rgb(hex_color):
    """Converts hex color string to RGBColor object."""
    if not hex_color or hex_color.lower() == 'black':
        return RGBColor(0, 0, 0)
    
    hex_color = hex_color.lstrip('#')
    hex_color = hex_color.split(' ')[0]
    
    try:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except ValueError:
        return RGBColor(0, 0, 0)

def apply_formatting(run, formatting, color_hex):
    """Applies formatting (bold, italic, size, color) to a text run."""
    if not formatting:
        return

    formatting = formatting.lower()
    
    if "bold" in formatting:
        run.font.bold = True
    if "italic" in formatting:
        run.font.italic = True
    
    if "large" in formatting:
        run.font.size = Pt(44)
    elif "medium" in formatting:
        run.font.size = Pt(32)
    elif "normal" in formatting:
        run.font.size = Pt(24)
    
    if color_hex:
        run.font.color.rgb = hex_to_rgb(color_hex)

def generate_slides(slides_data, output_file):
    """Generates slides from a dictionary of slide data."""
    prs = Presentation()
    
    for slide_num in sorted(slides_data.keys()):
        rows = slides_data[slide_num]
        slide_title_text = rows[0]['Slide Title']
        
        if slide_num == 1:
            layout = prs.slide_layouts[0] # Title Slide
        else:
            layout = prs.slide_layouts[1] # Title and Content
            
        slide = prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_title_text
            
        if slide_num == 1:
            for row in rows:
                element = row['Element']
                content = row['Content']
                formatting = row['Formatting']
                color = row['Color']
                
                if element == "Main Text":
                    shape = slide.shapes.title
                elif element == "Subtitle":
                    shape = slide.placeholders[1]
                else:
                    continue
                
                shape.text = content
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    if p.runs:
                        apply_formatting(p.runs[0], formatting, color)
                    else:
                        run = p.add_run()
                        run.text = content
                        apply_formatting(run, formatting, color)
        else:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for row in rows:
                element = row['Element']
                content = row['Content']
                formatting = row['Formatting']
                color = row['Color']
                
                if element == "Heading" and content == slide_title_text:
                    continue
                
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = content
                apply_formatting(run, formatting, color)
                p.space_after = Pt(10)

    prs.save(output_file)
    print(f"Successfully created {output_file}")

def create_presentation(csv_file, output_file):
    """Legacy wrapper for CSV files."""
    slides_data = {}
    try:
        with open(csv_file, 'r', encoding='utf-8-sig') as f:
            reader = csv.DictReader(f)
            for row in reader:
                slide_num = int(row['Slide Number'])
                if slide_num not in slides_data:
                    slides_data[slide_num] = []
                slides_data[slide_num].append(row)
    except FileNotFoundError:
        print(f"Error: Could not find file {csv_file}")
        return

    generate_slides(slides_data, output_file)

    generate_slides(slides_data, output_file)

def add_word_slides(prs, word_data, start_slide_num=1):
    """Adds slides for a single word to an existing presentation object."""
    # We need to map our "slides_data" structure to the actual presentation creation
    # The existing generate_slides function creates a NEW presentation.
    # We need to adapt the logic to append to an EXISTING one.
    
    # Helper to create slide data structure for one word
    slides_config = {}
    
    # Slide 1: Title
    slides_config[1] = [
        {"Slide Title": "Title Slide", "Element": "Main Text", "Content": f"Word of the Week: {word_data['word']}", "Formatting": "Large, Bold", "Color": "#186433"},
        {"Slide Title": "Title Slide", "Element": "Subtitle", "Content": "Spelling Focus", "Formatting": "Medium", "Color": "#186433"}
    ]
    
    # Slide 2: Definition
    slides_config[2] = [
        {"Slide Title": "Definition", "Element": "Heading", "Content": "Definition", "Formatting": "Large", "Color": "Black"},
        {"Slide Title": "Definition", "Element": "Content", "Content": word_data['definition'], "Formatting": "Normal", "Color": "Black"}
    ]
    
    # Slide 3: Sentence
    slides_config[3] = [
        {"Slide Title": "Usage in a Sentence", "Element": "Heading", "Content": "Usage in a Sentence", "Formatting": "Large", "Color": "Black"},
        {"Slide Title": "Usage in a Sentence", "Element": "Content", "Content": word_data['sentence'], "Formatting": "Normal", "Color": "Black"}
    ]
    
    # Slide 4: Word Origin & Parts
    origin_content = word_data.get('morphology') or word_data.get('etymology')
    if origin_content:
        slides_config[4] = [
            {"Slide Title": "Word Origin & Parts", "Element": "Heading", "Content": "Word Origin & Parts", "Formatting": "Large", "Color": "Black"},
            {"Slide Title": "Word Origin & Parts", "Element": "Content", "Content": origin_content, "Formatting": "Normal", "Color": "Black"}
        ]
        
    # Slide 5: Synonyms
    if word_data.get('synonyms'):
        synonyms_list = word_data['synonyms'].split(',')
        slide_rows = [{"Slide Title": "Synonyms", "Element": "Heading", "Content": "Synonyms", "Formatting": "Large", "Color": "Black"}]
        for syn in synonyms_list:
            slide_rows.append({"Slide Title": "Synonyms", "Element": "Content", "Content": syn.strip(), "Formatting": "Normal", "Color": "#186433"})
        slides_config[5] = slide_rows

    # Slide 6: Antonyms
    if word_data.get('antonyms'):
        antonyms_list = word_data['antonyms'].split(',')
        slide_rows = [{"Slide Title": "Antonyms", "Element": "Heading", "Content": "Antonyms", "Formatting": "Large", "Color": "Black"}]
        for ant in antonyms_list:
            slide_rows.append({"Slide Title": "Antonyms", "Element": "Content", "Content": ant.strip(), "Formatting": "Normal", "Color": "#B91C1C"})
        slides_config[6] = slide_rows

    # Now actually create the slides in the presentation
    for slide_key in sorted(slides_config.keys()):
        rows = slides_config[slide_key]
        slide_title_text = rows[0]['Slide Title']
        
        # Determine layout
        # If it's the very first slide of the word (Title Slide), use Title Layout
        if slide_key == 1:
            layout = prs.slide_layouts[0] # Title Slide
        else:
            layout = prs.slide_layouts[1] # Title and Content
            
        slide = prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_title_text
            
        if slide_key == 1:
            for row in rows:
                element = row['Element']
                content = row['Content']
                formatting = row['Formatting']
                color = row['Color']
                
                if element == "Main Text":
                    shape = slide.shapes.title
                elif element == "Subtitle":
                    shape = slide.placeholders[1]
                else:
                    continue
                
                shape.text = content
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    if p.runs:
                        apply_formatting(p.runs[0], formatting, color)
                    else:
                        run = p.add_run()
                        run.text = content
                        apply_formatting(run, formatting, color)
        else:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for row in rows:
                element = row['Element']
                content = row['Content']
                formatting = row['Formatting']
                color = row['Color']
                
                if element == "Heading" and content == slide_title_text:
                    continue
                
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = content
                apply_formatting(run, formatting, color)
                p.space_after = Pt(10)

def create_presentation_from_data(word_data, output_file):
    """Creates presentation from direct word data."""
    prs = Presentation()
    add_word_slides(prs, word_data)
    prs.save(output_file)
    print(f"Successfully created {output_file}")

def create_batch_presentation(list_of_word_data, output_file):
    """Creates a single presentation containing slides for multiple words."""
    prs = Presentation()
    for word_data in list_of_word_data:
        add_word_slides(prs, word_data)
    prs.save(output_file)
    print(f"Successfully created batch presentation {output_file}")

if __name__ == "__main__":
    create_presentation("Week-7-Spelling-Update-Guide.csv", "Spelling_Presentation.pptx")
